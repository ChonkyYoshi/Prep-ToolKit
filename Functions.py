import docx
import pptx
from regex import search, match, findall
from pathlib import Path
from win32com.client import DispatchEx
from zipfile import ZipFile
from PIL import Image
from shutil import copyfile, rmtree, make_archive, move
import helper as hp
import imagequant
from lxml import etree
import openpyxl as xl


def Upsave(File: Path):
    match File.suffix:
        case '.doc':
            File = Doc2Docx(File)
        case '.ppt':
            File = Ppt2Pptx(File)
        case '.xls':
            File = Xls2Xlsx(File)
    return File


def Doc2Docx(File: Path):

    FileStr = File.__str__()
    WordApp = DispatchEx('Word.Application')
    Doc = WordApp.Documents.Open(FileStr)
    Doc.SaveAs(FileStr, FileFormat=12)
    WordApp.Quit()
    File = Path(FileStr + '.docx')
    return File


def Xls2Xlsx(File: Path):

    XlApp = DispatchEx('Excel.Application')
    FileStr = File.__str__()
    Xl = XlApp.Workbooks.Open(FileStr)
    Xl.SaveAs(FileStr, FileFormat=51)
    XlApp.Quit()
    File = Path(FileStr + '.xlsx')
    return File


def Ppt2Pptx(File: Path):

    PptApp = DispatchEx('PowerPoint.Application')
    FileStr = File.__str__()
    Ppt = PptApp.Presentations.Open(FileStr, 0, 0, 0)
    Ppt.SaveAs(FileStr, FileFormat=24)
    PptApp.Quit()
    File = Path(FileStr + '.pptx')
    return File


def ExtractImages(File: Path):

    zipfile = ZipFile(File)
    TempDir = Path(File.parent.as_posix() + '/Temp')
    TempDir.mkdir(exist_ok=True)
    for file in zipfile.namelist():
        if match(r'.*?/media/.*?\.(jpeg|jpg|png|emf|wmf|wdp)', file):
            zipfile.extract(file, TempDir)
    if File.suffix == '.pptx' or File.suffix == '.story':
        for rel in zipfile.namelist():
            if match(r'.*?/slides/_rels', rel):
                zipfile.extract(rel, TempDir)
    return TempDir


def CleanTempDir(Tempdir: Path, compress=True):

    for index, file in enumerate(list(Tempdir.glob('*/media/*'))):
        i = len(list(Tempdir.glob('*/media/*')))
        yield f'Converting all to png, image {index + 1} of {i}'
        if not file.suffix == '.png':
            with Image.open(file.as_posix()) as im:
                try:
                    im.save(file.as_posix(), format='PNG')
                    file.unlink(missing_ok=True)
                except OSError:
                    continue
    if compress:
        for index, file in enumerate(list(Tempdir.glob('*/media/*'))):
            i = len(list(Tempdir.glob('*/media/*')))
            yield f'Compressing, image {index + 1} of {i}'
            if file.suffix == '.png':
                with Image.open(file.as_posix()) as im:
                    new_im = imagequant.quantize_pil_image(im)
                    new_im.save(file.as_posix()[:-3] + 'png')


def FillCS(TempDir: Path, File: Path):

    CSPath = Path(File.parent.as_posix() + '/Contact Sheets')
    CSPath.mkdir(exist_ok=True)
    CS = docx.Document()

    i = len(list(TempDir.glob('*/media/*')))
    for index, pic in enumerate(list(TempDir.glob('*/media/*'))):
        if pic.is_dir() or pic.suffix == '.rels':
            continue
        Table = CS.add_table(rows=5, cols=2, style='Table Grid')
        Table.cell(0, 0).merge(Table.cell(0, 1)).text = pic.name
        Table.cell(1, 0).merge(Table.cell(1, 1))
        Table.cell(2, 0).merge(Table.cell(2, 1))
        Table.cell(3, 0).text = 'Source'
        Table.cell(3, 1).text = 'Target'
        try:
            yield f'Filling in Contact Sheet, image {index + 1} of {i}'
            run = Table.cell(1, 0).paragraphs[0].add_run()
            run.add_picture(pic.as_posix(),
                            width=(CS.sections[0].page_width -
                                   (CS.sections[0].right_margin +
                                    CS.sections[0].left_margin)))
        except Exception:
            run = Table.cell(1, 0).paragraphs[0].add_run()
            run.text = f'''An error occured while trying to insert
the image, please check the Error folder and manually insert the image.
Name of the image: {pic.name}'''
            ErrorDir = Path(File.parent.as_posix() +
                            '/Errors_' + File.stem)
            ErrorDir.mkdir(exist_ok=True)
            copyfile(pic.as_posix(),
                     ErrorDir.as_posix() + '/' + pic.name)
        if File.suffix == '.pptx' or File.suffix == '.story':
            locations = LocateImage(TempDir, str(pic.name))
            if len(locations) == 0:
                Table.cell(2, 0).add_paragraph(
                    'only present in Master Slide', style='List Bullet')
                CS.add_page_break()
            else:
                for location in locations:
                    Table.cell(2, 0).add_paragraph(location,
                                                   style='List Bullet')
                CS.add_page_break()
        else:
            CS.add_page_break()
    CS.save(f'{CSPath.as_posix()}/CS_{File.name}.docx')
    rmtree(TempDir)


def LocateImage(TempDir: Path, ImageName: str):

    Locations = []
    for entry in TempDir.rglob('*.rels'):
        with open(entry, 'r') as File:
            rel = File.read()
        if search(ImageName, rel):
            match = search(r'(\w+)\.', entry.name)
            Locations.append(entry.
                             name[match.start():match.end()])  # type: ignore
    return Locations


def BilTable(File: Path):

    li = list()
    doc = docx.Document(File)
    for index, table in enumerate(doc.tables):
        yield f'processing table {index + 1} of {len(doc.tables)}'
        for c in range(len(table.columns)):
            for r in range(len(table.rows)):
                if table.cell(r, c)._tc not in li:
                    li.append(table.cell(r, c)._tc)
                    for par in table.cell(r, c).paragraphs:
                        prevpar = par.insert_paragraph_before()
                        hp.CopyParFormatting(prevpar, par)
                        for run in par.runs:
                            prevrun = prevpar.add_run()
                            prevrun.text = run.text
                            hp.CopyRunFormatting(prevrun, run)
                            prevrun.font.hidden = True
        li.clear()
    i = len(doc.paragraphs)
    for index, par in enumerate(doc.paragraphs):
        yield f'processing paragraph {index + 1} of {i}'
        table = doc.add_table(rows=1, cols=2)
        par._p.addnext(table._tbl)
        SPar = table.cell(0, 0).paragraphs[0]
        TPar = table.cell(0, 1).paragraphs[0]
        hp.CopyParFormatting(SPar, par)
        hp.CopyParFormatting(TPar, par)
        for run in par.runs:
            SRun = SPar.add_run()
            TRun = TPar.add_run()
            SRun.text = run.text
            TRun.text = run.text
            hp.CopyRunFormatting(SRun, run)
            hp.CopyRunFormatting(TRun, run)
        par._element.getparent().os.remove(par._element)
    doc.save(File.parent.as_posix() + '/Bil_' + File.name)


def Doc2PDF(WordApp, File: Path,
            ARev: bool = False,
            DRev: bool = False,
            Com: bool = False,
            Overwrite: bool = False):

    Doc = WordApp.Documents.Open(File.as_posix(), Visible=False)
    PdfDir = Path(File.as_posix() + '/PDF')
    PdfDir.mkdir(exist_ok=True)
    if Doc.Revisions.Count > 0 and ARev:
        Doc.AcceptAllRevisions()
    if Doc.Revisions.Count > 0 and DRev:
        Doc.RejectAllRevisions()
    if Doc.Comments.Count > 0 and Com:
        Doc.DeleteAllComments()
    if Overwrite:
        Doc.Save()
    Doc.SaveAs2(f'{PdfDir.as_posix()}/{File.name}.pdf', FileFormat=17)
    Doc.Close()


def AcceptRevisions(WordApp, File: Path,
                    ARev: bool = False,
                    DRev: bool = False,
                    Com: bool = False,
                    Overwrite: bool = False):

    Doc = WordApp.Documents.Open(File.as_posix(), Visible=False)

    if Doc.Revisions.Count > 0 and ARev:
        Doc.AcceptAllRevisions()
    if Doc.Revisions.Count > 0 and DRev:
        Doc.RejectAllRevisions()
    if Doc.Comments.Count > 0 and Com:
        Doc.DeleteAllComments()
    if Overwrite:
        Doc.Save()
    else:
        match File.suffix:
            case '.docx':
                Doc.SaveAs2(File.parent.as_posix() + 'NoRev_' +
                            File.stem, FileFormat=12)
            case '.docm':
                Doc.SaveAs2(File.parent.as_posix() + 'NoRev_' +
                            File.stem, FileFormat=13)
            case '.doc':
                Doc.SaveAs2(File.parent.as_posix() + 'NoRev_' +
                            File.stem, FileFormat=0)
    Doc.Close()


def PrepStoryExport(File: Path, Regex: str = ''):

    Doc = docx.Document(File)
    for index, par in enumerate(Doc.paragraphs):
        yield f'paragraph {index + 1} of {len(Doc.paragraphs)}'
        for run in par.runs:
            run.font.hidden = True
    for table in Doc.tables:
        for index, SkipCol in enumerate(table.columns):
            yield f'column {index + 1} of {len(table.columns)}'
            if not index == 3:
                for cell in SkipCol.cells:
                    for par in cell.paragraphs:
                        for run in par.runs:
                            run.font.hidden = True
            else:
                for cell in SkipCol.cells:
                    for par in cell.paragraphs:
                        for run in par.runs:
                            if match(Regex, run.text) and Regex != '':
                                start = match(Regex,
                                              run.text).start()  # type: ignore
                                end = match(Regex,
                                            run.text).end()  # type: ignore
                                hidden_run = hp.isolate_run(par, start, end)
                                hidden_run.font.hidden = True
                for par in table.cell(0, 3).paragraphs:
                    for run in par.runs:
                        run.font.hidden = True
    Doc.save(File.parent.as_posix() + '/Prep_' + File.name)


def Unhide(File: Path,
           SkipRow: bool = False,
           SkipCol: bool = False,
           SkipSheet: bool = False,
           SkipShp: bool = False,
           SkipSld: bool = False,
           Overwrite: bool = False):

    match File.suffix:
        case '.docx' | '.docm':
            Doc = docx.Document(File)
            for index, par in enumerate(Doc.paragraphs):
                yield f'Paragraph {index + 1} of {len(Doc.Paragraphs)}'
                for run in par.runs:
                    run.font.hidden = False
            for index, table in enumerate(Doc.tables):
                yield f'Table {index + 1} of {len(Doc.tables)}'
                for row in table.rows:
                    for cell in row.cells:
                        for par in cell.paragraphs:
                            for run in par.runs:
                                run.font.hidden = False
            for section in Doc.sections:
                for par in section.header.paragraphs:
                    for run in par.runs:
                        run.font.hidden = False
                for par in section.footer.paragraphs:
                    for run in par.runs:
                        run.font.hidden = False
            if Overwrite:
                Doc.save(File.as_posix())
            else:
                Doc.save(f'{File.parent.as_posix()}/UNH_{File.name}')
        case '.xlsx' | '.xlsm':
            Sheets = list()
            wb = xl.load_workbook(File, read_only=True)
            for ws in wb.worksheets:
                if ws.sheet_state != 'visible':
                    Sheets.append(f'sheet{wb.get_index(ws) + 1}.xml')
            wb.close()
            Temp = Path(f'{File.parent.as_posix()}/Temp')
            Path.mkdir(Temp, exist_ok=True)
            ZipFile(File).extractall(Temp)
            for Sheetindex, Sheetfile in enumerate(
                    list(Temp.rglob('xl/worksheets/*.xml'))):
                Sheetcount = len(list(Temp.rglob('xl/worksheets/*.xml')))
                Sheetfile = Path(Sheetfile)
                Sheetxml = etree.parse(source=Sheetfile,
                                       parser=etree.XMLParser())
                if Sheetfile.name in Sheets and SkipSheet:
                    continue
                else:
                    wbfile = Path(f'{Temp.as_posix()}/xl/workbook.xml')
                    wbxml = etree.parse(source=wbfile,
                                        parser=etree.XMLParser())
                    for sheet in wbxml.xpath('//*[local-name()="sheet"]'):
                        sheet.set("state", "visible")
                if not SkipRow:
                    RowCount = len(list(
                        Sheetxml.xpath('//*[local-name()="row"]')))
                    for Rowindex, row in enumerate(
                            Sheetxml.xpath('//*[local-name()="row"]')):
                        yield f'sheet {Sheetindex + 1} of {Sheetcount}' +\
                            f'\nrow {Rowindex + 1} of {RowCount}'
                        row.set("hidden", "0")
                if not SkipCol:
                    ColCount = len(list(
                        Sheetxml.xpath('//*[local-name()="col"]')))
                    for ColIndex, col in enumerate(
                            Sheetxml.xpath('//*[local-name()="col"]')):
                        yield f'sheet {Sheetindex + 1} of {Sheetcount}' +\
                            f'\ncolumn {ColIndex + 1} of {ColCount}'
                        col.set('hidden', '0')
                with open(Sheetfile, 'wb') as f:
                    f.write(etree.tostring(Sheetxml))
            if Overwrite:
                new = Path(make_archive(File.stem, 'zip', root_dir=Temp))
                move(new, File.as_posix())
            else:
                new = Path(make_archive(f'UNH_{File.stem}',
                                        'zip', root_dir=Temp))
                move(new, f'{File.parent.as_posix()}/{new.stem}{File.suffix}')
            rmtree(Temp)
        case '.pptx' | '.pptm':
            Pres = pptx.Presentation(File)
            for index, slide in enumerate(Pres.slides):
                yield f'Slide {index} of {len(Pres.slides)}'
                if not SkipSld:
                    slide._element.set('show', '1')
                if not SkipShp:
                    for shape in slide.shapes:
                        hide(shape)
            if Overwrite:
                Pres.save(File)
            else:
                Pres.save(f'{File.parent.as_posix()}/UNH_{File.name}')


def hide(shape):

    match shape.shape_type:
        case 6:
            for sub in shape.shapes:
                hide(sub)
        case 13:
            shape._element.nvPicPr.cNvPr.set('hidden', '0')
        case 9:
            shape._element.nvCxnSpPr.cNvPr.set('hidden', '0')
        case 3 | 19 | 7:
            shape._element.nvGraphicFramePr.cNvPr.set('hidden', '0')
        case _:
            try:
                shape._element.nvSpPr.cNvPr.set('hidden', '0')
            except Exception:
                breakpoint()


def PPTSections(File: Path):

    PPT = pptx.Presentation(File)
    xml = str(PPT.part.blob)
    if search(r'(<p14:section name=")(.*?)(")', xml):
        doc = docx.Document()
        table = doc.add_table(rows=1, cols=2)
        r = table.cell(0, 0).paragraphs[0].add_run()
        r.text = 'Source'
        r.font.hidden = True
        r = table.cell(0, 1).paragraphs[0].add_run()
        r.text = 'Target'
        r.font.hidden = True
        count = len(findall(r'(<p14:section name=")(.*?)(")', xml))
        for index, sec in enumerate(findall(r'(<p14:section name=")(.*?)(")',
                                            xml)):
            yield f'section {index + 1} of {count}'
            row = table.add_row()
            r = row.cells[0].paragraphs[0].add_run()
            r.text = sec[1]
            r.font.hidden = True
            r = row.cells[1].paragraphs[0].add_run()
            r.text = sec[1]
        doc.save(f'{File.parent.as_posix()}/Section Titles_{File.name}.docx')
